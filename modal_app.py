"""
modal_app.py — Resonance Analyzer Extraction Service on Modal
Replaces Flask + ngrok with a permanent cloud endpoint.
Deploy with: modal deploy modal_app.py
"""

import modal
from pathlib import Path

app = modal.App("resonance-analyzer-extractor")

image = (
    modal.Image.debian_slim(python_version="3.11")
    .pip_install([
        "fastapi[standard]",
        "flask",
        "pdfplumber",
        "python-docx",
        "python-pptx",
        "openpyxl",
        "Pillow",
        "pytesseract",
        "requests",
    ])
    .apt_install(["tesseract-ocr"])
)


@app.function(
    image=image,
    timeout=300,
)
@modal.fastapi_endpoint(method="POST")
def extract(item: dict):
    import tempfile
    import os
    import base64
    import csv
    import requests as req
    from pathlib import Path
    import pdfplumber
    from docx import Document
    from pptx import Presentation
    from openpyxl import load_workbook
    from PIL import Image
    import pytesseract

    def extract_text_from_pdf(file_path):
        text = []
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text.append(page_text)
        return "\n\n".join(text)

    def extract_text_from_docx(file_path):
        doc = Document(file_path)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

    def extract_text_from_pptx(file_path):
        prs = Presentation(file_path)
        text = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"--- Slide {slide_num} ---"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            text.append("\n".join(slide_text))
        return "\n\n".join(text)

    def extract_text_from_xlsx(file_path):
        workbook = load_workbook(file_path, data_only=True)
        text = []
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            text.append(f"--- Sheet: {sheet_name} ---")
            for row in sheet.iter_rows(values_only=True):
                row_text = [str(cell) for cell in row if cell is not None]
                if row_text:
                    text.append("\t".join(row_text))
        return "\n".join(text)

    def extract_text_from_csv(file_path):
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            reader = csv.reader(f)
            rows = list(reader)
        if not rows:
            return ""
        return "\n".join(",".join(row) for row in rows if any(cell.strip() for cell in row))

    def extract_text_from_image(file_path):
        image = Image.open(file_path)
        return pytesseract.image_to_string(image).strip()

    def extract_text(file_path, filename):
        ext = Path(filename).suffix.lower()
        extractors = {
            '.pdf': extract_text_from_pdf,
            '.docx': extract_text_from_docx,
            '.pptx': extract_text_from_pptx,
            '.xlsx': extract_text_from_xlsx,
            '.csv': extract_text_from_csv,
            '.png': extract_text_from_image,
            '.jpg': extract_text_from_image,
            '.jpeg': extract_text_from_image,
        }
        extractor = extractors.get(ext)
        if not extractor:
            raise Exception(f"Unsupported file format: {ext}")
        return extractor(file_path)

    files = item.get("files", [])
    results = []

    for file_info in files:
        filename = file_info.get("filename", "unknown")
        result = {"filename": filename}
        temp_path = None

        try:
            suffix = Path(filename).suffix
            with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
                temp_path = tmp.name

            if "content_b64" in file_info:
                file_bytes = base64.b64decode(file_info["content_b64"])
                with open(temp_path, 'wb') as f:
                    f.write(file_bytes)
            elif "url" in file_info:
                response = req.get(file_info["url"], timeout=60, stream=True)
                response.raise_for_status()
                with open(temp_path, 'wb') as f:
                    for chunk in response.iter_content(chunk_size=8192):
                        f.write(chunk)
            else:
                raise Exception("No file content provided")

            result["extracted_text"] = extract_text(temp_path, filename)

        except Exception as e:
            result["error"] = str(e)

        finally:
            if temp_path and os.path.exists(temp_path):
                os.unlink(temp_path)

        results.append(result)

    return {"files": results}


if __name__ == "__main__":
    print("Deploy with: modal deploy modal_app.py")
    print("Test with: modal run modal_app.py")