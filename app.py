from flask import Flask, request, jsonify
import os
import tempfile
import requests
from pathlib import Path
import pdfplumber
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
from PIL import Image
import pytesseract

app = Flask(__name__)


def download_file(url, filename):
    """Download a file from URL to a temporary location."""
    try:
        response = requests.get(url, timeout=30, stream=True)
        response.raise_for_status()

        # Create temporary file with proper extension
        suffix = Path(filename).suffix
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=suffix)

        with open(temp_file.name, 'wb') as f:
            for chunk in response.iter_content(chunk_size=8192):
                f.write(chunk)

        return temp_file.name
    except Exception as e:
        raise Exception(f"Failed to download file: {str(e)}")


def extract_text_from_pdf(file_path):
    """Extract text from PDF using pdfplumber."""
    text = []
    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text.append(page_text)
    return "\n\n".join(text)


def extract_text_from_docx(file_path):
    """Extract text from DOCX using python-docx."""
    doc = Document(file_path)
    text = []
    for paragraph in doc.paragraphs:
        if paragraph.text.strip():
            text.append(paragraph.text)
    return "\n".join(text)


def extract_text_from_pptx(file_path):
    """Extract text from PPTX using python-pptx."""
    prs = Presentation(file_path)
    text = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text = [f"--- Slide {slide_num} ---"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text)
        text.append("\n".join(slide_text))

    return "\n\n".join(text)


def extract_text_from_xlsx(file_path):
    """Extract text from XLSX using openpyxl."""
    workbook = load_workbook(file_path, data_only=True)
    text = []

    for sheet_name in workbook.sheetnames:
        sheet = workbook[sheet_name]
        text.append(f"--- Sheet: {sheet_name} ---")

        for row in sheet.iter_rows(values_only=True):
            row_text = [str(cell) if cell is not None else "" for cell in row]
            row_text = [cell for cell in row_text if cell.strip()]
            if row_text:
                text.append("\t".join(row_text))

    return "\n".join(text)

def extract_text_from_csv(file_path):
    """Extract text from CSV files."""
    import csv
    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
        reader = csv.reader(f)
        rows = list(reader)
    if not rows:
        return ""
    lines = [",".join(row) for row in rows if any(cell.strip() for cell in row)]
    return "\n".join(lines)

def extract_text_from_image(file_path):
    """Extract text from image using Tesseract OCR."""
    image = Image.open(file_path)
    text = pytesseract.image_to_string(image)
    return text.strip()


def extract_text(file_path, filename):
    """Extract text based on file extension."""
    extension = Path(filename).suffix.lower()

    extractors = {
        '.pdf': extract_text_from_pdf,
        '.docx': extract_text_from_docx,
        '.pptx': extract_text_from_pptx,
        '.xlsx': extract_text_from_xlsx,
        '.csv': extract_text_from_csv,
        '.png': extract_text_from_image,
        '.jpg': extract_text_from_image,
        '.jpeg': extract_text_from_image,
    }

    extractor = extractors.get(extension)
    if not extractor:
        raise Exception(f"Unsupported file format: {extension}")

    return extractor(file_path)


@app.route('/extract', methods=['POST'])
def extract():
    """
    Extract text from multiple files.

    Expected JSON body:
    {
        "files": [
            {"url": "https://example.com/file.pdf", "filename": "document.pdf"},
            {"url": "https://example.com/image.png", "filename": "image.png"}
        ]
    }

    Returns:
    {
        "files": [
            {"filename": "document.pdf", "extracted_text": "..."},
            {"filename": "image.png", "error": "..."}
        ]
    }
    """
    try:
        data = request.get_json()

        if not data or 'files' not in data:
            return jsonify({"error": "Missing 'files' array in request body"}), 400

        if not isinstance(data['files'], list):
            return jsonify({"error": "'files' must be an array"}), 400

        results = []

        for file_info in data['files']:
            result = {"filename": file_info.get('filename', 'unknown')}
            temp_path = None

            try:
                # Validate file info
                if 'url' not in file_info or 'filename' not in file_info:
                    raise Exception("Each file must have 'url' and 'filename' fields")

                url = file_info['url']
                filename = file_info['filename']

                # Download file
                temp_path = download_file(url, filename)

                # Extract text
                extracted_text = extract_text(temp_path, filename)
                result['extracted_text'] = extracted_text

            except Exception as e:
                result['error'] = str(e)

            finally:
                # Clean up temporary file
                if temp_path and os.path.exists(temp_path):
                    try:
                        os.unlink(temp_path)
                    except:
                        pass

            results.append(result)

        return jsonify({"files": results}), 200

    except Exception as e:
        return jsonify({"error": f"Request failed: {str(e)}"}), 500


@app.route('/', methods=['GET'])
def index():
    """Health check endpoint."""
    return jsonify({
        "service": "Text Extraction API",
        "status": "running",
        "supported_formats": ["PDF", "DOCX", "PPTX", "XLSX", "PNG", "JPG", "JPEG"]
    }), 200


if __name__ == '__main__':
    app.run(debug=True, host='0.0.0.0', port=5000)
